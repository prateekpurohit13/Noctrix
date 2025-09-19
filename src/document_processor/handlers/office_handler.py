import io
from pathlib import Path
from typing import List, Tuple
import docx
import openpyxl
from pptx import Presentation
from PIL import Image

from src.common.models import Element, TextElement, TableElement, TableCell, ImageElement
from src.document_processor.ocr_service import perform_ocr

def _extract_images_from_docx(doc: docx.Document, temp_dir: Path) -> List[Element]:
    image_elements = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image_bytes = rel.target_part.blob
                image_path = temp_dir / Path(rel.target_ref).name
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                
                print(f"  -> Found embedded image in DOCX. OCRing '{image_path.name}'...")
                ocr_text = perform_ocr(image_path)
                image_elements.append(ImageElement(type="image", text_content=ocr_text, image_path=str(image_path)))
            except Exception as e:
                print(f"  -> WARNING: Could not extract image from DOCX. Reason: {e}")
    return image_elements

def _extract_images_from_pptx(prs: Presentation, temp_dir: Path) -> List[Element]:
    image_elements = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext
                    image_path = temp_dir / f"slide{i+1}_img.{ext}"
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)

                    print(f"  -> Found embedded image in PPTX. OCRing '{image_path.name}'...")
                    ocr_text = perform_ocr(image_path)
                    image_elements.append(ImageElement(type="image", text_content=ocr_text, image_path=str(image_path)))
                except Exception as e:
                    print(f"  -> WARNING: Could not extract image from PPTX. Reason: {e}")
    return image_elements
    
def _extract_images_from_xlsx(workbook: openpyxl.Workbook, temp_dir: Path) -> List[Element]:
    image_elements = []
    for sheet in workbook.worksheets:
        for i, image in enumerate(sheet._images):
            try:
                image_bytes = image.ref
                pil_image = Image.open(io.BytesIO(image_bytes))
                ext = pil_image.format.lower()
                image_path = temp_dir / f"{sheet.title}_img{i+1}.{ext}"
                pil_image.save(image_path)
                
                print(f"  -> Found embedded image in XLSX. OCRing '{image_path.name}'...")
                ocr_text = perform_ocr(image_path)
                image_elements.append(ImageElement(type="image", text_content=ocr_text, image_path=str(image_path)))
            except Exception as e:
                print(f"  -> WARNING: Could not extract image from XLSX. Reason: {e}")
    return image_elements


def process_office_doc(file_path: Path, temp_dir: Path) -> Tuple[List[Element], int]:
    print(f"  -> Processing as Office Doc: '{file_path.name}'")
    extension = file_path.suffix.lower()
    elements: List[Element] = []
    page_count = 1

    if extension == ".docx":
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                elements.append(TextElement(type="paragraph", text_content=para.text.strip()))
        elements.extend(_extract_images_from_docx(doc, temp_dir))

    elif extension == ".xlsx":
        workbook = openpyxl.load_workbook(file_path)
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            rows_data = [[TableCell(row=i, col=j, text_content=str(cell.value)) for j, cell in enumerate(row) if cell.value is not None] for i, row in enumerate(sheet.iter_rows())]
            elements.append(TableElement(type="table", rows=rows_data, metadata={'sheet_name': sheet_name}))
        elements.extend(_extract_images_from_xlsx(workbook, temp_dir))

    elif extension == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    elements.append(TextElement(type="paragraph", text_content=shape.text.strip()))
        elements.extend(_extract_images_from_pptx(prs, temp_dir))
        page_count = len(prs.slides)
        
    return elements, page_count